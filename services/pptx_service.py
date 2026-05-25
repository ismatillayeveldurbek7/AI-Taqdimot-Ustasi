"""
AI Taqdimot Ustasi — Python-pptx bilan PPTX generator
Node.js kerak emas, faqat: pip install python-pptx Pillow aiohttp
"""

import io
import asyncio
import aiohttp
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

THEMES = {
    "Blue": {
        "bg":      RGBColor(0x0D, 0x1B, 0x3E),
        "title":   RGBColor(0x00, 0xC9, 0xFF),
        "text":    RGBColor(0xFF, 0xFF, 0xFF),
        "accent":  RGBColor(0x4A, 0x90, 0xD9),
        "card":    RGBColor(0x1A, 0x37, 0x6C),
        "point":   RGBColor(0xA8, 0xC8, 0xF0),
    },
    "Black": {
        "bg":      RGBColor(0x1C, 0x1C, 0x1C),
        "title":   RGBColor(0xFF, 0xD7, 0x00),
        "text":    RGBColor(0xFF, 0xFF, 0xFF),
        "accent":  RGBColor(0xE0, 0xA8, 0x00),
        "card":    RGBColor(0x33, 0x33, 0x33),
        "point":   RGBColor(0xAA, 0xAA, 0xAA),
    },
    "White": {
        "bg":      RGBColor(0xFF, 0xFF, 0xFF),
        "title":   RGBColor(0x2E, 0x86, 0xAB),
        "text":    RGBColor(0x1C, 0x1C, 0x1C),
        "accent":  RGBColor(0xA2, 0x3B, 0x72),
        "card":    RGBColor(0xF0, 0xF4, 0xFF),
        "point":   RGBColor(0x55, 0x55, 0x55),
    },
    "Green": {
        "bg":      RGBColor(0x10, 0x47, 0x2F),
        "title":   RGBColor(0x7E, 0xD3, 0x21),
        "text":    RGBColor(0xFF, 0xFF, 0xFF),
        "accent":  RGBColor(0x2E, 0xCC, 0x71),
        "card":    RGBColor(0x1A, 0x6B, 0x3C),
        "point":   RGBColor(0xA8, 0xE6, 0xBC),
    },
    "PremiumDark": {
        "bg":      RGBColor(0x0A, 0x0A, 0x0F),
        "title":   RGBColor(0xC0, 0x8C, 0xFF),
        "text":    RGBColor(0xF0, 0xF0, 0xFF),
        "accent":  RGBColor(0x7B, 0x2F, 0xFF),
        "card":    RGBColor(0x1A, 0x0A, 0x2E),
        "point":   RGBColor(0xD0, 0xB8, 0xFF),
    },
}

W = Inches(13.33)
H = Inches(7.5)


def _rgb(color: RGBColor):
    return color


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text(slide, text, x, y, w, h, font_size, color: RGBColor,
              bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


async def _download_image(url: str) -> bytes | None:
    try:
        async with aiohttp.ClientSession() as sess:
            async with sess.get(url, timeout=aiohttp.ClientTimeout(total=8)) as r:
                if r.status == 200:
                    ct = r.headers.get("Content-Type", "")
                    if "image" in ct:
                        return await r.read()
    except Exception:
        pass
    return None


def _add_cover_slide(prs, title: str, theme: dict):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, theme["bg"])

    # Accent bar top
    _add_rect(slide, 0, 0, W, Inches(0.08), theme["accent"])

    # Title
    _add_text(
        slide, title,
        Inches(0.7), Inches(2.2),
        Inches(11.9), Inches(2.2),
        font_size=40, color=theme["title"],
        bold=True, align=PP_ALIGN.CENTER
    )

    # Subtitle
    _add_text(
        slide, "AI Taqdimot Ustasi tomonidan yaratildi",
        Inches(0.7), Inches(4.6),
        Inches(11.9), Inches(0.6),
        font_size=16, color=theme["point"],
        align=PP_ALIGN.CENTER
    )

    # Accent bar bottom
    _add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), theme["accent"])
    return slide


def _add_content_slide(prs, slide_data: dict, theme: dict, img_bytes: bytes | None = None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, theme["bg"])

    # Top accent
    _add_rect(slide, 0, 0, W, Inches(0.06), theme["accent"])

    # Slide number badge
    num = slide_data.get("number", "")
    _add_rect(slide, Inches(0.3), Inches(0.2), Inches(0.55), Inches(0.42), theme["accent"])
    _add_text(slide, str(num),
              Inches(0.3), Inches(0.2), Inches(0.55), Inches(0.42),
              font_size=14, color=theme["bg"], bold=True, align=PP_ALIGN.CENTER)

    # Title
    slide_title = slide_data.get("title", "")
    _add_text(slide, slide_title,
              Inches(1.0), Inches(0.18),
              Inches(11.0), Inches(0.6),
              font_size=22, color=theme["title"],
              bold=True, align=PP_ALIGN.LEFT)

    # Separator line
    _add_rect(slide, Inches(0.3), Inches(0.88), Inches(12.7), Inches(0.03), theme["accent"])

    has_image = img_bytes is not None
    content_width = Inches(7.5) if has_image else Inches(12.5)

    # Content text
    content = slide_data.get("content", "")
    if content:
        _add_text(slide, content,
                  Inches(0.4), Inches(1.05),
                  content_width, Inches(2.0),
                  font_size=14, color=theme["text"],
                  align=PP_ALIGN.LEFT)

    # Key points
    key_points = slide_data.get("key_points", [])
    if key_points:
        _add_text(slide, "Asosiy fikrlar:",
                  Inches(0.4), Inches(3.15),
                  content_width, Inches(0.38),
                  font_size=13, color=theme["accent"],
                  bold=True, align=PP_ALIGN.LEFT)

        for i, point in enumerate(key_points[:5]):
            y = Inches(3.55) + i * Inches(0.62)
            # bullet dot
            _add_rect(slide, Inches(0.4), y + Inches(0.12),
                      Inches(0.12), Inches(0.12), theme["accent"])
            _add_text(slide, point,
                      Inches(0.62), y,
                      content_width - Inches(0.3), Inches(0.58),
                      font_size=13, color=theme["point"],
                      align=PP_ALIGN.LEFT)

    # Image (right side)
    if has_image:
        try:
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(
                img_stream,
                Inches(8.1), Inches(1.05),
                Inches(4.9), Inches(5.5)
            )
        except Exception:
            pass

    # Bottom bar
    _add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), theme["accent"])
    return slide


async def generate_pptx(data: dict, color_scheme: str = "Blue") -> bytes:
    theme = THEMES.get(color_scheme, THEMES["Blue"])
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_data = data.get("slides", [])

    # Download images concurrently
    img_tasks = []
    for sd in slides_data:
        url = sd.get("image_url")
        if url:
            img_tasks.append(_download_image(url))
        else:
            img_tasks.append(asyncio.sleep(0, result=None))

    images = await asyncio.gather(*img_tasks)

    # Cover slide
    _add_cover_slide(prs, data.get("title", "Taqdimot"), theme)

    # Content slides
    for sd, img in zip(slides_data, images):
        _add_content_slide(prs, sd, theme, img)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
